"""
Document Processing Module for various document types
"""
import os
import json
from pathlib import Path
from typing import Dict, Any, List, Optional
from dataclasses import dataclass, asdict

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import PyPDF2
import markdown
from bs4 import BeautifulSoup
from loguru import logger


@dataclass
class DocumentChunk:
    """Represents a chunk of document content"""
    content: str
    metadata: Dict[str, Any]
    page_number: Optional[int] = None
    slide_number: Optional[int] = None


class DocumentProcessor:
    """Process various document formats (PDF, DOCX, PPTX, HTML, etc.)"""
    
    def __init__(self, chunk_size: int = 500, overlap: int = 50):
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def process_document(self, file_path: Path, metadata: Optional[Dict] = None) -> List[DocumentChunk]:
        """
        Process document based on file extension
        
        Args:
            file_path: Path to document
            metadata: Additional metadata
        
        Returns:
            List of document chunks
        """
        ext = file_path.suffix.lower()
        metadata = metadata or {}
        metadata['source_file'] = str(file_path)
        
        if ext == '.pdf':
            return self._process_pdf(file_path, metadata)
        elif ext == '.docx':
            return self._process_docx(file_path, metadata)
        elif ext in ['.pptx', '.ppt']:
            return self._process_pptx(file_path, metadata)
        elif ext == '.md':
            return self._process_markdown(file_path, metadata)
        elif ext in ['.html', '.htm']:
            return self._process_html(file_path, metadata)
        elif ext in ['.txt', '.json', '.csv']:
            return self._process_text(file_path, metadata)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return []
    
    def _process_pdf(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process PDF file"""
        chunks = []
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        # Split into chunks if needed
                        chunk_metadata = metadata.copy()
                        chunk_metadata['page'] = page_num + 1
                        chunk_metadata['document_type'] = 'pdf'
                        
                        # Create chunks from page content
                        page_chunks = self._split_text(text, chunk_metadata)
                        chunks.extend(page_chunks)
                        
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
        
        return chunks
    
    def _process_docx(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process DOCX file"""
        chunks = []
        try:
            doc = DocxDocument(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            text = '\n'.join(full_text)
            chunk_metadata = metadata.copy()
            chunk_metadata['document_type'] = 'docx'
            
            chunks = self._split_text(text, chunk_metadata)
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
        
        return chunks
    
    def _process_pptx(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process PPTX file"""
        chunks = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)
                
                if slide_text:
                    text = ' '.join(slide_text)
                    chunk_metadata = metadata.copy()
                    chunk_metadata['slide'] = slide_num + 1
                    chunk_metadata['document_type'] = 'pptx'
                    chunks.append(
                        DocumentChunk(
                            content=text,
                            metadata=chunk_metadata,
                            slide_number=slide_num + 1
                        )
                    )
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
        
        return chunks
    
    def _process_markdown(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Convert markdown to text
            html = markdown.markdown(content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            chunk_metadata = metadata.copy()
            chunk_metadata['document_type'] = 'markdown'
            
            return self._split_text(text, chunk_metadata)
            
        except Exception as e:
            logger.error(f"Error processing Markdown {file_path}: {e}")
            return []
    
    def _process_html(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process HTML file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
            
            text = soup.get_text()
            chunk_metadata = metadata.copy()
            chunk_metadata['document_type'] = 'html'
            
            return self._split_text(text, chunk_metadata)
            
        except Exception as e:
            logger.error(f"Error processing HTML {file_path}: {e}")
            return []
    
    def _process_text(self, file_path: Path, metadata: Dict) -> List[DocumentChunk]:
        """Process plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            chunk_metadata = metadata.copy()
            chunk_metadata['document_type'] = 'text'
            
            return self._split_text(text, chunk_metadata)
            
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            return []
    
    def _split_text(self, text: str, metadata: Dict) -> List[DocumentChunk]:
        """Split text into chunks with overlap"""
        if not text.strip():
            return []
        
        # Simple sentence-based chunking
        sentences = text.split('. ')
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_len = len(sentence)
            if current_length + sentence_len > self.chunk_size and current_chunk:
                chunks.append(
                    DocumentChunk(
                        content='. '.join(current_chunk) + '.',
                        metadata=metadata.copy()
                    )
                )
                # Keep overlap
                overlap_text = current_chunk[-1] if current_chunk else ''
                current_chunk = [overlap_text] if overlap_text else []
                current_length = len(overlap_text)
            
            current_chunk.append(sentence)
            current_length += sentence_len
        
        if current_chunk:
            chunks.append(
                DocumentChunk(
                    content='. '.join(current_chunk) + '.',
                    metadata=metadata.copy()
                )
            )
        
        return chunks